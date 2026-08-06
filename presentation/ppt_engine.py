import calendar
import math
import re
from copy import deepcopy
from pathlib import Path
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.oxml import serialize_part_xml
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from pptx.text.text import _Paragraph
from pptx.util import Emu, Pt

from presentation.slide_ops import (
    add_table_column, delete_slide, duplicate_slide, set_column_widths
)

# Stopgap font sizes for objects where the template run is styled as a
# short heading/placeholder but currently receives a full AI-generated
# paragraph. Remove once each object gets its own properly-designed
# multi-paragraph AI mapping.
FONT_SIZE_OVERRIDES = {
    "AI_ExecutiveSummary": Pt(12),
}

# Body font for the fixed-size "label + 3 lines" commentary cards. The
# template sized these for its own example sentences; real AI commentary
# citing this client's figures runs longer, and these cards are narrow enough
# (~22 characters a line) that 14pt text wraps to 4 lines per bullet and pushes
# the last one out through the bottom of the card. Dropping the body a couple
# of points buys a line per bullet, which is what keeps full sentences intact
# instead of truncating them mid-clause.
COMMENTARY_BODY_FONT = {
    "AI_TrendAnalysisSummary": Pt(11),
    "AI_EngagementSummary": Pt(11),
    "AI_ContentPerformanceSummary": Pt(11),
}

# Vertical space at the foot of every slide occupied by the template's
# confidentiality line and logo. Footnote boxes are parked just above this
# rather than relative to the content above them -- see fit_footnote().
FOOTER_BAND_HEIGHT = Emu(640000)

# Explicit per-series colors for the comparison chart. The template's own
# Q1/Q2 colors lead the list so the 2-period case looks unchanged;
# replace_data() otherwise reuses the last series' exact color for every
# series added beyond the original two, making 3+ periods visually
# indistinguishable.
COMPARISON_SERIES_COLORS = [
    "0F3D63", "1C6BFF", "2FA0A0", "F2A93B",
    "8E6FCE", "5A6B7B", "E0607E", "4C9F70",
]

# StatusPill_* shapes carry a fixed fill color per slide position in the
# template (e.g. whichever metric row happened to be "Strong Growth" in
# its own example data got the green pill). Real data can put any status
# at any row, so the pill's color needs to track its own text instead of
# staying wherever the template originally painted it. Colors taken from
# the template's own existing pills (Decline has no example in the
# template's own data, so it reuses the theme's accent2 -- the same
# orange already used for the "Status" header pill).
STATUS_PILL_COLORS = {
    "Strong Growth": "1EA86B",
    "Stable Growth": "1C6AFE",
    "No Change": "9A9FB3",
    "Decline": "E97132",
}

class PowerPointEngine:

    def __init__(self):

        self.template = Path("templates/LeadGen_QBR_Template.pptx")
        self.prs = None

    # ---------------------------------------------------------
    # Load PowerPoint
    # ---------------------------------------------------------

    def load(self):

        self.prs = Presentation(self.template)

    # ---------------------------------------------------------
    # Find object by Selection Pane name
    # (supports grouped shapes)
    # ---------------------------------------------------------

    def find_shape(self, object_name, occurrence=0):

        """
        Returns the `occurrence`-th shape (0-indexed) named `object_name`
        across the whole deck, in slide order. Most objects only appear
        once, so the default finds it as before. Some names are legitimately
        reused across slides (e.g. PERIOD_Q1 appears on both the Campaign
        Snapshot slide and the Q1 detail slide) -- pass occurrence=1 to
        reach the second one.
        """

        def search(shapes):

            for shape in shapes:

                if shape.name == object_name:
                    yield shape

                # Search inside grouped shapes
                if hasattr(shape, "shapes"):
                    yield from search(shape.shapes)

        seen = 0

        for slide in self.prs.slides:

            for shape in search(slide.shapes):

                if seen == occurrence:
                    return shape

                seen += 1

        return None

    # ---------------------------------------------------------
    # Replace text
    # ---------------------------------------------------------

    def replace_text(self, object_name, value, occurrence=0, paragraph_index=0, group_child_index=None, run_index=None):

        shape = self.find_shape(object_name, occurrence=occurrence)

        if shape is None:

            print(f"[NOT FOUND] {object_name}")

            return False

        # -------------------------------------------------
        # Grouped KPI cards
        # -------------------------------------------------

        def write_into(tf, wrap=False):

            # The number boxes are sized (spAutoFit) for whatever digit count
            # the template's example happened to have. With word_wrap on, a
            # longer real value wraps to a second line that overlaps the rest
            # of the card instead of growing the box -- word_wrap off lets it
            # overflow horizontally instead, the safer failure mode.
            #
            # Captions are the opposite case. They are full phrases sitting in
            # a box the template already sized for two wrapped lines ("Total
            # Leads across 5 markets"), so switching wrapping off runs them
            # straight out through the side of the card and across the card
            # beside it. They keep whatever wrapping the template gave them.
            if not wrap:
                tf.word_wrap = False

            if tf.paragraphs and tf.paragraphs[0].runs:

                tf.paragraphs[0].runs[0].text = str(value)

            else:

                tf.clear()

                tf.paragraphs[0].text = str(value)

        if not shape.has_text_frame:

            if hasattr(shape, "shapes"):

                children = list(shape.shapes)

                # Some cards (e.g. the Country count) need a specific
                # child targeted -- the label below the number, not the
                # number itself -- rather than "whichever child has text
                # first", which always finds the number.
                if group_child_index is not None:

                    if (
                        group_child_index < len(children)
                        and children[group_child_index].has_text_frame
                    ):
                        # group_child_index only ever targets a card's caption
                        # -- the number is found by the search below.
                        write_into(
                            children[group_child_index].text_frame, wrap=True
                        )
                        print(f"[UPDATED GROUP] {object_name}")
                        return True

                else:

                    # KPI card groups include a background auto-shape
                    # that also reports has_text_frame (empty,
                    # decorative) -- target the child that already
                    # carries the placeholder number/text, not just the
                    # first text frame found.
                    for child in children:

                        if child.has_text_frame and child.text_frame.text.strip():

                            write_into(child.text_frame)
                            print(f"[UPDATED GROUP] {object_name}")
                            return True

            print(f"[NO TEXT FRAME] {object_name}")

            return False

        tf = shape.text_frame

        if paragraph_index >= len(tf.paragraphs):

            print(f"[NO SUCH PARAGRAPH] {object_name}[{paragraph_index}]")

            return False

        paragraph = tf.paragraphs[paragraph_index]

        # `is not None`, not a truthiness test: run_index=0 is a real target
        # (the figure run of a "1,032  Total Leads" pair) and must NOT fall
        # through to the branch below, which strips every run after the first
        # and would delete the label sitting in run 1.
        if run_index is not None:

            # Some paragraphs split a fixed prefix from their body across two
            # runs -- AI_H2Recommendations numbers each line with a bold
            # "01   " run followed by the text run. Writing run 0 there would
            # wipe the numbering, so target the body run and leave the rest
            # (and the run count) alone.
            if run_index < len(paragraph.runs):

                paragraph.runs[run_index].text = str(value)

            else:

                print(
                    f"[NO SUCH RUN] {object_name}"
                    f"[p{paragraph_index}][r{run_index}]"
                )

                return False

        elif paragraph.runs:

            # Preserve formatting by editing the first run
            paragraph.runs[0].text = str(value)

            # Remove any extra runs
            while len(paragraph.runs) > 1:
                paragraph._element.remove(paragraph.runs[-1]._r)

        else:

            paragraph.text = str(value)

        if (
            paragraph_index == 0
            and object_name in FONT_SIZE_OVERRIDES
            and tf.paragraphs[0].runs
        ):
            tf.paragraphs[0].runs[0].font.size = FONT_SIZE_OVERRIDES[object_name]

        # StatusPill_* shapes carry a fixed fill color per slide position
        # in the template -- recolor to match the actual status text so
        # a "Decline" pill doesn't stay whatever color the template's
        # own example happened to paint that row.
        if object_name.startswith("StatusPill_") and str(value) in STATUS_PILL_COLORS:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor.from_string(
                STATUS_PILL_COLORS[str(value)]
            )

        print(f"[UPDATED] {object_name}")

        return True

    def clone_paragraphs(self, object_name, source_indices):

        """
        Appends copies of existing paragraphs to a fixed-length text box, so a
        list the template drew with N entries can carry N+1.

        Key Learnings and the recommendations list are single text boxes whose
        entries are plain paragraphs, not list placeholders -- there is no way
        to add an entry except to add a paragraph. Copying an existing one
        carries its font, weight, colour and spacing across, which writing a
        fresh paragraph would not.

        Returns the index of the first appended paragraph, or None if the
        object isn't there.
        """

        shape = self.find_shape(object_name)

        if shape is None or not shape.has_text_frame:
            print(f"[SKIPPED] clone_paragraphs {object_name} (not found)")
            return None

        paragraphs = shape.text_frame.paragraphs

        first_new = len(paragraphs)

        sources = []

        for index in source_indices:

            if not -len(paragraphs) <= index < len(paragraphs):
                print(f"[SKIPPED] clone_paragraphs {object_name}[{index}]")
                return None

            sources.append(paragraphs[index]._p)

        # Resolved to elements before any is appended -- appending shifts the
        # negative indices the callers use, so reading them lazily would copy
        # a copy on the second pass.
        body = shape.text_frame._txBody

        for source in sources:
            body.append(deepcopy(source))

        print(f"[EXTENDED] {object_name} by {len(source_indices)} paragraph(s)")

        return first_new

    # ---------------------------------------------------------
    # Replace a variable-length numbered list (e.g. the agenda),
    # cloning the template's first paragraph so run-level
    # formatting (font, color, numbering prefix style) survives
    # regardless of how many items the selected mode produces.
    # ---------------------------------------------------------

    def replace_list(self, object_name, items):

        shape = self.find_shape(object_name)

        if shape is None:

            print(f"[NOT FOUND] {object_name}")

            return False

        if not shape.has_text_frame or not shape.text_frame.paragraphs:

            print(f"[NO TEXT FRAME] {object_name}")

            return False

        tf = shape.text_frame
        txBody = tf._txBody

        # The template's own example has a fixed number of items sized to
        # exactly fill the box at its own font size/spacing -- a mode
        # that produces more items (e.g. Month over Month with several
        # months) than that baseline would otherwise overflow the box
        # top and bottom, since nothing here scales font size or spacing
        # down to compensate. Scale both proportionally when there are
        # more items than the template's own baseline.
        baseline_item_count = len(tf.paragraphs)
        baseline_run = tf.paragraphs[0].runs[0] if tf.paragraphs[0].runs else None
        baseline_font_size = baseline_run.font.size if baseline_run is not None else None
        baseline_space_after = tf.paragraphs[0].space_after

        scale = (
            min(1.0, baseline_item_count / len(items))
            if items else 1.0
        )

        template_p = deepcopy(tf.paragraphs[0]._p)

        for paragraph in list(tf.paragraphs):
            txBody.remove(paragraph._p)

        for i, label in enumerate(items, start=1):

            new_p_element = deepcopy(template_p)
            txBody.append(new_p_element)

            paragraph = _Paragraph(new_p_element, tf)
            runs = paragraph.runs

            if len(runs) >= 2:
                runs[0].text = f"{i:02d}   "
                runs[1].text = label
            elif len(runs) == 1:
                runs[0].text = f"{i:02d}   {label}"

            if scale < 1.0:

                if baseline_font_size is not None:

                    new_size = Pt(round(baseline_font_size.pt * scale))

                    for run in runs:
                        run.font.size = new_size

                if baseline_space_after is not None:
                    paragraph.space_after = Pt(round(baseline_space_after.pt * scale))

        print(f"[UPDATED LIST] {object_name} ({len(items)} items, scale={scale:.2f})")

        return True

    def replace_table(self, object_name, dataframe):

        print(f"[TABLE] {object_name}")

        shape = self.find_shape(object_name)

        if shape is None:

            print(f"[NOT FOUND] {object_name}")

            return False

        if not shape.has_table:

            print(f"[NOT A TABLE] {object_name}")

            return False

        table = shape.table

        # A client's report can legitimately come back with zero rows
        # (e.g. no trending topics detected for that program/period).
        # Leaving the loop below to run zero times would silently keep
        # the template's own example rows on screen -- specific,
        # plausible-looking company names/scores that have nothing to
        # do with this client -- which reads as real data, not as
        # "nothing here". Blank every data row instead so it's visibly
        # empty rather than quietly wrong.
        if dataframe.empty:

            for r in range(1, len(table.rows)):
                for c in range(len(table.columns)):
                    self._write_cell(table.cell(r, c), "")

            print(f"[CLEARED TABLE] {object_name} (no data)")

            return True

        # ---------------------------------------
        # Write dataframe into PPT table
        # ---------------------------------------

        max_rows = min(len(dataframe), len(table.rows) - 1)

        max_cols = min(len(dataframe.columns), len(table.columns))

        for r in range(max_rows):

            for c in range(max_cols):

                value = dataframe.iloc[r, c]

                if value is None or (isinstance(value, float) and math.isnan(value)):

                    value = ""

                elif isinstance(value, bool):

                    pass

                elif isinstance(value, (int, float)):

                    # Every numeric table cell in this deck is either a
                    # whole-number count (leads, accounts, a Change
                    # delta) or a score that should read as one (e.g. an
                    # ML score averaging out to 96.66666) -- round and
                    # thousands-format uniformly instead of hunting down
                    # each table's own source. Values that need a "%"
                    # suffix (e.g. % Change) are formatted as strings
                    # before they ever reach this table, so they pass
                    # through untouched here.
                    value = f"{round(value):,}"

                self._write_cell(table.cell(r + 1, c), value)

        # A client's data can legitimately fill fewer rows than the
        # template's table has (e.g. only 5 accounts showed intent, but
        # the table was drawn with 8 data rows). The loop above only
        # touches the rows it had data for, which left the template's own
        # example rows -- real-looking company names and scores from
        # whatever deck this template was built from -- sitting
        # underneath the genuine data, indistinguishable from it. Blank
        # every row the data didn't reach.
        leftover = range(max_rows + 1, len(table.rows))

        for r in leftover:
            for c in range(len(table.columns)):
                self._write_cell(table.cell(r, c), "")

        if leftover:
            print(
                f"[UPDATED TABLE] {object_name} "
                f"({max_rows} rows, {len(leftover)} unused row(s) cleared)"
            )
        else:
            print(f"[UPDATED TABLE] {object_name}")

        return True

    @staticmethod
    def _write_cell(cell, value):

        """Writes a cell's text through its first existing run where
        there is one, so the template's own font/size/colour for that
        cell survives (assigning to text_frame.text drops it)."""

        tf = cell.text_frame

        if tf.paragraphs and tf.paragraphs[0].runs:
            tf.paragraphs[0].runs[0].text = str(value)
        else:
            tf.text = str(value)
    
    def replace_chart(self, object_name, dataframe):

        print(f"[CHART] {object_name}")

        shape = self.find_shape(object_name)

        if shape is None:

            print(f"[NOT FOUND] {object_name}")

            return False

        if not shape.has_chart:

            print(f"[NOT A CHART] {object_name}")

            return False

        chart = shape.chart

        # A client's report can legitimately come back with zero rows
        # (e.g. no trending topics detected for that program/period) --
        # python-pptx can't build a chart with no categories at all, and
        # there's nothing meaningful to plot anyway, so leave the
        # template's own chart in place rather than crashing the whole
        # generation run over one empty source table.
        if dataframe is None or dataframe.empty:

            print(f"[SKIPPED] {object_name} (no data)")

            return False

        # "Is Forecast" (Chart_TrendAnalysis) marks which trailing
        # points should render dashed -- it's a styling hint, not a
        # series to plot, so pull it out before building categories/series.
        is_forecast = None

        if "Is Forecast" in dataframe.columns:
            is_forecast = dataframe["Is Forecast"].tolist()
            dataframe = dataframe.drop(columns=["Is Forecast"])

        chart_data = CategoryChartData()

        # Categories = first column
        chart_data.categories = dataframe.iloc[:, 0].tolist()

        # Remaining columns become series
        for column in dataframe.columns[1:]:

            chart_data.add_series(
                str(column),
                dataframe[column].tolist()
            )

        chart.replace_data(chart_data)

        # Bar width scales down as more periods add more series per
        # category, so the comparison chart stays readable whether it's
        # showing 2 periods or (Month over Month) several.
        if object_name == "CHART_Q1Q2Comparison":

            num_series = len(dataframe.columns) - 1

            chart.plots[0].gap_width = max(30, 150 - 15 * num_series)

            self._recolor_series(chart, num_series)

            if num_series != 2:
                self._reset_legend_layout(chart)

        if object_name == "Chart_TrendAnalysis":

            # The template's value axis has a manual minimum (300) tuned
            # for its original example data -- real/forecast values for a
            # different campaign can fall well under or over that. Rather
            # than leaving the axis on pure auto-scale (which PowerPoint
            # usually starts near/at 0, making small real swings look
            # flat), start it at a "nice" rounded-down value just below
            # the actual minimum across every series (e.g. a minimum of
            # 755 starts the axis at 700) so the real variation reads
            # clearly.
            numeric_columns = dataframe.columns[1:]

            min_value = (
                dataframe[numeric_columns].min().min()
                if len(numeric_columns) else None
            )

            chart.value_axis.minimum_scale = (
                self._axis_floor(min_value) if min_value is not None else None
            )
            chart.value_axis.maximum_scale = None

            if is_forecast:
                self._style_forecast_points(chart, is_forecast)

        if object_name == "Chart_ContentPerformance":

            # The template highlights one bar (the top performer) in a
            # distinct color via a single per-point override, fixed at
            # whatever index happened to be the winner in the template's
            # own example data. Real data sorts differently per client,
            # so keep that highlight pointed at whichever bar is
            # actually highest, not a hardcoded position.
            value_column = dataframe.columns[1]
            self._highlight_top_n_points(chart, dataframe[value_column].tolist(), 1)

        if object_name == "Chart_TrendingTopics":

            # Same idea, but the template highlights 3 bars at fixed
            # positions (3, 8, 9) from its own example data. Highlight
            # whichever 2 bars are actually highest instead.
            value_column = dataframe.columns[1]
            self._highlight_top_n_points(chart, dataframe[value_column].tolist(), 2)

        if object_name == "Chart_TopicDistribution":

            self._style_topic_distribution(chart, dataframe)

        if object_name == "Chart_BuyingStage":

            # The template color-codes each Predictive Buying Stage bar
            # (grey = no signal, light blue = active/sales-ready, ML
            # teal = the default) via per-point overrides fixed at the
            # index order of its own example data. Real data is sorted
            # by account count, so a different dataset can put the
            # stages in a different order -- rebind each color to the
            # stage name wherever it actually lands instead of a
            # hardcoded position.
            self._style_buying_stage(chart, dataframe.iloc[:, 0].tolist())

        if object_name == "Chart_EngagementFunnel":

            # The template's single-series funnel colors each of its 4
            # bars individually (one dPt per stage) -- that scheme no
            # longer fits now that the funnel can carry a second
            # "All Accounts" series alongside "Trending". Use one color
            # per series instead (same approach as the comparison chart).
            num_series = len(dataframe.columns) - 1

            self._recolor_series(chart, num_series)

            # Legend shown even for a single series. When the campaign has no
            # Target Account List History export there is no "All Accounts"
            # series to plot, and suppressing the legend left the chart with
            # nothing saying the bars cover trending accounts only -- easily
            # read as the whole account universe, and materially lower than the
            # platform's own funnel for the same campaign.
            chart.has_legend = True

            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False

        # Comma-format every count shown directly on the chart itself
        # (data labels, value axis) so a 4+ digit number reads "1,234"
        # like every other count in the deck, not "1234". Topic
        # Distribution is the one chart whose labels are percentages,
        # not counts, and already sets its own "0%" format -- skip it
        # here so this doesn't clobber that.
        if object_name != "Chart_TopicDistribution":
            self._format_labels_with_commas(chart)

        print(f"[UPDATED CHART] {object_name}")

        return True

    @staticmethod
    def _format_labels_with_commas(chart):

        for plot in chart.plots:

            try:
                plot.data_labels.number_format = "#,##0"
                plot.data_labels.number_format_is_linked = False
            except (ValueError, AttributeError):
                pass

        try:
            chart.value_axis.tick_labels.number_format = "#,##0"
            chart.value_axis.tick_labels.number_format_is_linked = False
        except (ValueError, AttributeError):
            pass

    @staticmethod
    def _axis_floor(value):

        """Rounds a value down to a "nice" step sized to its own
        magnitude (nearest lower 100 for a value in the hundreds,
        nearest lower 1000 for a value in the thousands, etc.) -- e.g.
        755 -> 700, 3371 -> 3000. Used to give the trend chart's value
        axis a minimum close to the real data instead of starting at 0
        or an unrelated fixed number."""

        if value is None or value <= 0:
            return 0

        step = 10 ** math.floor(math.log10(value))

        return math.floor(value / step) * step

    @staticmethod
    def _recolor_series(chart, num_series):

        """replace_data() reuses the last existing series' exact color
        for every series added beyond the template's original two, so
        3+ periods render as visually indistinguishable bars. Assign an
        explicit, distinct color to each series instead.

        The template also carries a per-point color override (<c:dPt>)
        on the first category for the original Q2 series, which
        replace_data() copies onto every series cloned from it -- that
        override wins over the series-level fill set here, so the
        cloned series' first bar (Total Leads) kept rendering in the
        old override color while its other bars picked up the new one.
        Dropping any <c:dPt> elements lets the series-level fill apply
        uniformly across every category.
        """

        for series, color in zip(chart.plots[0].series, COMPARISON_SERIES_COLORS):

            series.format.fill.solid()
            series.format.fill.fore_color.rgb = RGBColor.from_string(color)

            for dPt in series._element.findall(qn("c:dPt")):
                series._element.remove(dPt)

    @staticmethod
    def _reset_legend_layout(chart):

        """The template's legend has a manual layout box sized for two
        short labels (Q1/Q2), positioned to the right of the bars. With
        3+ (longer) period labels, entries overflow that fixed box and
        silently get clipped from view -- removing the manual layout
        lets PowerPoint size it to fit, and moving it above the plot
        (rather than squeezed to the right) gives it room to lay out
        several entries in a row instead of clipping again."""

        if not chart.has_legend:
            return

        legend_elm = chart._chartSpace.find(qn("c:chart")).find(qn("c:legend"))

        if legend_elm is None:
            return

        layout_elm = legend_elm.find(qn("c:layout"))

        if layout_elm is not None:
            legend_elm.remove(layout_elm)

        chart.legend.position = XL_LEGEND_POSITION.TOP
        chart.legend.include_in_layout = False

    @staticmethod
    def _style_forecast_points(chart, is_forecast):

        """Renders the trailing "Is Forecast" points dashed, like the
        template's original Q3/Q4 styling -- but for however many
        forecast points there actually are (2-3, depending on mode),
        not just the two the template happened to be built with.

        The template already carries one correctly-styled <c:dPt> per
        series (dashed line + dashed marker outline) for exactly this
        purpose; reusing it as a clone source guarantees valid schema
        placement (dPt must sit between marker and dLbls) instead of
        hand-building the XML from scratch.
        """

        num_real = is_forecast.index(True) if True in is_forecast else len(is_forecast)
        num_forecast = len(is_forecast) - num_real

        if num_forecast <= 0:
            return

        for series in chart.plots[0].series:

            series_elm = series._element

            existing_dPts = series_elm.findall(qn("c:dPt"))

            if not existing_dPts:
                continue

            template_dPt = deepcopy(existing_dPts[0])

            for dPt in existing_dPts:
                series_elm.remove(dPt)

            anchor = series_elm.find(qn("c:marker"))

            if anchor is None:
                anchor = series_elm.find(qn("c:spPr"))

            for offset in range(num_forecast):

                new_dPt = deepcopy(template_dPt)
                new_dPt.find(qn("c:idx")).set("val", str(num_real + offset))

                anchor.addnext(new_dPt)
                anchor = new_dPt

    def _highlight_top_n_points(self, chart, values, n):

        """Keeps the template's highlight-color data point(s) (a
        distinct fill on the top performer bar(s)) pointed at whichever
        bars actually have the N highest values, rather than the fixed
        indices baked in from the template's own example data (which
        only matched its real top performers by coincidence of how
        that example happened to be sorted)."""

        if not values:
            return

        series_list = list(chart.plots[0].series)

        if not series_list:
            return

        series_elm = series_list[0]._element
        existing_dPts = series_elm.findall(qn("c:dPt"))

        if not existing_dPts:
            return

        template_dPt = deepcopy(existing_dPts[0])

        for dPt in existing_dPts:
            series_elm.remove(dPt)

        ranked = sorted(range(len(values)), key=lambda i: values[i], reverse=True)
        top_indices = sorted(ranked[:n])

        # Note: lxml elements are falsy when childless, so this must use
        # explicit None checks -- an `or` chain would skip a valid but
        # childless match (e.g. <c:invertIfNegative val="0"/>).
        anchor = series_elm.find(qn("c:invertIfNegative"))

        if anchor is None:
            anchor = series_elm.find(qn("c:marker"))

        if anchor is None:
            anchor = series_elm.find(qn("c:spPr"))

        for idx in top_indices:

            new_dPt = deepcopy(template_dPt)
            new_dPt.find(qn("c:idx")).set("val", str(idx))

            # Some of the template's own highlighted bars only tint the
            # outline and leave the bar's own fill at the series default
            # (navy, same as every other bar) -- a border alone reads as
            # "no visible highlight" at a glance. Force the actual fill
            # to the highlight color too, so every highlighted bar looks
            # the same regardless of which chart's dPt it was cloned from.
            fill_color = new_dPt.find(f"{qn('c:spPr')}/{qn('a:solidFill')}")

            if fill_color is not None:

                for child in list(fill_color):
                    fill_color.remove(child)

                fill_color.append(
                    parse_xml(f'<a:srgbClr xmlns:a="{self.DRAWING_NS}" val="1C6BFF"/>')
                )

            anchor.addnext(new_dPt)
            anchor = new_dPt

    def _style_buying_stage(self, chart, categories):

        """Rebind the Predictive Buying Stage bar colors to each stage's
        name rather than its position: grey for "No Active Signals",
        light blue for the active/sales-ready stages, and the series'
        own ML teal default for everything else (Awareness today, any
        future stage tomorrow)."""

        GREY_STAGES = {"No Active Signals"}
        BLUE_STAGES = {"Decision", "Consideration", "Preawareness"}

        series_list = list(chart.plots[0].series)

        if not series_list:
            return

        series_elm = series_list[0]._element
        existing_dPts = series_elm.findall(qn("c:dPt"))

        if not existing_dPts:
            return

        template_dPt = deepcopy(existing_dPts[0])

        for dPt in existing_dPts:
            series_elm.remove(dPt)

        # Note: lxml elements are falsy when childless, so this must use
        # explicit None checks -- an `or` chain would skip a valid but
        # childless match.
        anchor = series_elm.find(qn("c:invertIfNegative"))

        if anchor is None:
            anchor = series_elm.find(qn("c:marker"))

        if anchor is None:
            anchor = series_elm.find(qn("c:spPr"))

        for idx, category in enumerate(categories):

            if category in GREY_STAGES:
                fill_xml = (
                    f'<a:schemeClr xmlns:a="{self.DRAWING_NS}" val="bg2">'
                    f'<a:lumMod val="75000"/></a:schemeClr>'
                )
            elif category in BLUE_STAGES:
                fill_xml = f'<a:srgbClr xmlns:a="{self.DRAWING_NS}" val="1C6BFF"/>'
            else:
                # Leave at the series-level default (ML teal) -- no
                # per-point override needed.
                continue

            new_dPt = deepcopy(template_dPt)
            new_dPt.find(qn("c:idx")).set("val", str(idx))

            fill_color = new_dPt.find(f"{qn('c:spPr')}/{qn('a:solidFill')}")

            for child in list(fill_color):
                fill_color.remove(child)

            fill_color.append(parse_xml(fill_xml))

            anchor.addnext(new_dPt)
            anchor = new_dPt

        # The template's category axis label rotation is an odd value
        # (-60000000, i.e. -1000 degrees -- OOXML angles wrap modulo
        # 360, so this actually renders as a steep ~80 degree slant, not
        # the small tilt it looks like at a glance) that reads as
        # slanted/near-vertical rather than horizontal. Force it flat.
        cat_ax_elm = chart._chartSpace.find(
            f"{{{self.CHART_NS}}}chart/{{{self.CHART_NS}}}plotArea/{{{self.CHART_NS}}}catAx"
        )

        if cat_ax_elm is not None:

            body_pr = cat_ax_elm.find(
                f"{{{self.CHART_NS}}}txPr/{{{self.DRAWING_NS}}}bodyPr"
            )

            if body_pr is not None:
                body_pr.set("rot", "0")

            # A stage taxonomy with more categories than the template's
            # own example (e.g. a client using a custom stage the
            # standard 5 don't cover) can exceed however many labels
            # PowerPoint decides fit horizontally, so it auto-skips every
            # other one rather than overlapping them -- silently hiding
            # real category names. Force every label to show regardless.
            existing_skip = cat_ax_elm.find(f"{{{self.CHART_NS}}}tickLblSkip")

            if existing_skip is not None:
                existing_skip.set("val", "1")
            else:

                skip_elm = parse_xml(
                    f'<c:tickLblSkip xmlns:c="{self.CHART_NS}" val="1"/>'
                )

                lbl_offset = cat_ax_elm.find(f"{{{self.CHART_NS}}}lblOffset")

                if lbl_offset is not None:
                    lbl_offset.addnext(skip_elm)
                else:
                    cat_ax_elm.append(skip_elm)

            # More categories than the template's own baseline (5)
            # means less horizontal room per label -- shrink the font
            # proportionally so labels are more likely to wrap cleanly
            # within their own slot instead of spilling into neighbors.
            baseline_categories = 5

            if len(categories) > baseline_categories:

                def_rpr = cat_ax_elm.find(
                    f"{{{self.CHART_NS}}}txPr/{{{self.DRAWING_NS}}}p"
                    f"/{{{self.DRAWING_NS}}}pPr/{{{self.DRAWING_NS}}}defRPr"
                )

                if def_rpr is not None and def_rpr.get("sz"):

                    baseline_size = int(def_rpr.get("sz"))

                    scale = baseline_categories / len(categories)

                    def_rpr.set("sz", str(round(baseline_size * scale)))

    # ---------------------------------------------------------
    # Topic Distribution doughnut: per-slice color, %-formatted bold
    # labels with contrast-aware text color, small slices' labels
    # nudged outward, and the real total in the chart's center overlay.
    # ---------------------------------------------------------

    CHART_NS = "http://schemas.openxmlformats.org/drawingml/2006/chart"
    DRAWING_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

    @staticmethod
    def _contrast_text_color(hex_color):

        """Perceived-luminance check: returns a dark or light hex color,
        whichever reads clearly against the given background."""

        r = int(hex_color[0:2], 16) / 255
        g = int(hex_color[2:4], 16) / 255
        b = int(hex_color[4:6], 16) / 255

        luminance = 0.299 * r + 0.587 * g + 0.114 * b

        return "0F3D63" if luminance > 0.6 else "FFFFFF"

    @staticmethod
    def _radial_offset(start_fraction, end_fraction, magnitude=0.09, angle_nudge_degrees=0):

        """Outward (x, y) nudge for a data label, based on the slice's
        angular midpoint -- doughnut/pie charts have no native "outside"
        label position (unlike bar/column charts), so a small slice's
        label has to be manually pushed clear of the ring via a layout
        offset, same mechanism the template already used by hand for
        its own small slices.

        `angle_nudge_degrees` rotates the push direction away from the
        slice's true angle -- needed when several small slices sit only
        a few degrees apart (e.g. four slices of 2-3% each clustered in
        a ~30 degree arc): pushing each one further out along its own
        near-identical angle still leaves their labels overlapping, since
        radial lines that close together stay close together regardless
        of distance. Fanning the angle itself, in addition to the
        distance, is what actually separates them."""

        midpoint = (start_fraction + end_fraction) / 2
        angle = math.radians(midpoint * 360 + angle_nudge_degrees)

        return magnitude * math.sin(angle), -magnitude * math.cos(angle)

    def _style_topic_distribution(self, chart, dataframe):

        values = dataframe.iloc[:, 1].tolist()
        total = sum(values) or 1

        series_list = list(chart.plots[0].series)

        if not series_list:
            return

        series_elm = series_list[0]._element

        # ---- per-slice colors ----

        for dPt in series_elm.findall(qn("c:dPt")):
            series_elm.remove(dPt)

        colors = [
            COMPARISON_SERIES_COLORS[i % len(COMPARISON_SERIES_COLORS)]
            for i in range(len(values))
        ]

        anchor = series_elm.find(qn("c:tx"))

        for i, color in enumerate(colors):

            dPt_elm = parse_xml(
                f'<c:dPt xmlns:c="{self.CHART_NS}" xmlns:a="{self.DRAWING_NS}">'
                f'<c:idx val="{i}"/>'
                f'<c:bubble3D val="0"/>'
                f'<c:spPr>'
                f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
                f'<a:ln w="19050"><a:solidFill><a:schemeClr val="lt1"/></a:solidFill></a:ln>'
                f'</c:spPr>'
                f"</c:dPt>"
            )

            anchor.addnext(dPt_elm)
            anchor = dPt_elm

        # ---- per-slice data labels: %, bold, contrast color, and an
        # outward nudge for slices too small to hold readable text ----

        dLbls_elm = series_elm.find(qn("c:dLbls"))

        if dLbls_elm is not None:
            series_elm.remove(dLbls_elm)

        small_slice_threshold = 0.06
        cumulative = 0.0
        small_slice_rank = 0

        dLbl_xml_parts = []

        for i, (color, value) in enumerate(zip(colors, values)):

            share = value / total
            start_fraction = cumulative
            cumulative += share

            text_color = self._contrast_text_color(color)

            layout_xml = ""

            if share < small_slice_threshold:

                # A pure radial push isn't enough on its own: several
                # small slices next to each other sit at nearly the same
                # angle, so pushing them all out by the same amount just
                # moves the collision outward instead of resolving it --
                # and even pushing them different *distances* isn't
                # enough when the angles themselves are only a few
                # degrees apart (radial lines that close together stay
                # close together regardless of length). Each successive
                # small slice (in angular order) gets pushed both further
                # out AND at a wider angle off its own true position,
                # alternating left/right of it, so the labels genuinely
                # fan apart instead of clustering along nearly-parallel
                # lines.
                magnitude = 0.14 + small_slice_rank * 0.06

                # Always rotate the same direction (rather than
                # alternating) so successive small slices fan out
                # monotonically in their natural angular order, with no
                # risk of two labels' push directions crossing back over
                # each other.
                angle_nudge = small_slice_rank * 9

                small_slice_rank += 1

                dx, dy = self._radial_offset(
                    start_fraction, cumulative, magnitude, angle_nudge
                )

                layout_xml = (
                    "<c:layout><c:manualLayout>"
                    f'<c:x val="{dx}"/><c:y val="{dy}"/>'
                    "</c:manualLayout></c:layout>"
                )

                # Contrast was chosen against the slice's own fill, but a
                # pushed-out label sits on the slide's light background
                # instead, not the slice -- so the "opposite of a dark
                # slice" white text becomes invisible on a white page.
                text_color = "0F3D63"

            dLbl_xml_parts.append(
                f'<c:dLbl xmlns:c="{self.CHART_NS}" xmlns:a="{self.DRAWING_NS}">'
                f'<c:idx val="{i}"/>'
                f"{layout_xml}"
                f"<c:spPr><a:noFill/><a:ln><a:noFill/></a:ln></c:spPr>"
                f"<c:txPr><a:bodyPr/><a:lstStyle/><a:p><a:pPr>"
                f'<a:defRPr sz="1197" b="1">'
                f'<a:solidFill><a:srgbClr val="{text_color}"/></a:solidFill>'
                f"</a:defRPr></a:pPr><a:endParaRPr lang=\"en-US\"/></a:p></c:txPr>"
                f'<c:showLegendKey val="0"/><c:showVal val="0"/>'
                f'<c:showCatName val="0"/><c:showSerName val="0"/>'
                f'<c:showPercent val="1"/><c:showBubbleSize val="0"/>'
                f"</c:dLbl>"
            )

        dLbls_xml = (
            f'<c:dLbls xmlns:c="{self.CHART_NS}" xmlns:a="{self.DRAWING_NS}">'
            + "".join(dLbl_xml_parts)
            + '<c:numFmt formatCode="0%" sourceLinked="0"/>'
            + "<c:spPr><a:noFill/><a:ln><a:noFill/></a:ln></c:spPr>"
            + '<c:showLegendKey val="0"/><c:showVal val="0"/>'
            + '<c:showCatName val="0"/><c:showSerName val="0"/>'
            + '<c:showPercent val="1"/><c:showBubbleSize val="0"/>'
            + "</c:dLbls>"
        )

        new_dLbls_elm = parse_xml(dLbls_xml)

        cat_elm = series_elm.find(qn("c:cat"))
        cat_elm.addprevious(new_dLbls_elm)

        # ---- fix the ring's own position, so its center is known ----
        #
        # The plot area had no manual layout, so PowerPoint auto-sizes it
        # to whatever's left after the legend's own fixed-width manual
        # layout -- normally fine, but the "N Mentions" text box is
        # anchored at a fixed fraction of the *whole chart* (ring + legend
        # combined) carried over from the template's own example. Real
        # category labels are longer than that example's, and legend
        # text wrapping can shift how much of the chart the auto-sized
        # ring ends up occupying, so the fixed text box increasingly sits
        # off-center from the ring's actual hole. Giving the ring itself
        # a fixed layout removes that guesswork -- its center becomes a
        # known quantity the text box can be aligned to exactly.
        legend_x = 0.6

        legend_elm = chart._chartSpace.find(f"{{{self.CHART_NS}}}chart/{{{self.CHART_NS}}}legend")

        if legend_elm is not None:

            legend_x_elm = legend_elm.find(
                f"{{{self.CHART_NS}}}layout/{{{self.CHART_NS}}}manualLayout/{{{self.CHART_NS}}}x"
            )

            if legend_x_elm is not None:
                legend_x = float(legend_x_elm.get("val"))

        plot_area_elm = chart._chartSpace.find(f"{{{self.CHART_NS}}}chart/{{{self.CHART_NS}}}plotArea")

        old_plot_layout = plot_area_elm.find(qn("c:layout"))

        if old_plot_layout is not None:
            plot_area_elm.remove(old_plot_layout)

        plot_width = max(legend_x - 0.02, 0.1)

        new_plot_layout = parse_xml(
            f'<c:layout xmlns:c="{self.CHART_NS}">'
            "<c:manualLayout>"
            '<c:xMode val="edge"/><c:yMode val="edge"/>'
            f'<c:x val="0"/><c:y val="0"/>'
            f'<c:w val="{plot_width}"/><c:h val="1"/>'
            "</c:manualLayout></c:layout>"
        )

        plot_area_elm.insert(0, new_plot_layout)

        # ---- center overlay: real total, not the template's example ----

        self._update_center_total(chart, total, ring_center=(plot_width / 2, 0.5))

    CHART_DRAWING_NS = "http://schemas.openxmlformats.org/drawingml/2006/chartDrawing"

    def _update_center_total(self, chart, total, ring_center=None):

        """The "N Mentions" text in the doughnut's hole is a chart
        user-shape (a text box anchored to the chart, not a normal
        slide shape or chart element) -- python-pptx has no API for
        this, so it's read/patched as a plain XML part."""

        try:
            user_shapes_part = chart.part.part_related_by(RT.CHART_USER_SHAPES)
        except KeyError:
            return

        root = parse_xml(user_shapes_part.blob)

        number_run = root.find(f".//{{{self.DRAWING_NS}}}p/{{{self.DRAWING_NS}}}r/{{{self.DRAWING_NS}}}t")

        if number_run is not None:
            number_run.text = f"{total:,.0f}"

            # The template's text box isn't vertically centered (no
            # anchor="ctr" on bodyPr, so it defaults to top-aligned) --
            # unnoticeable with its own short example number, but a
            # longer real total (more digits, or a comma) shifts the
            # 2-line "N / Mentions" block so it visibly hugs the top of
            # the ring instead of sitting centered in the hole.
            body_pr = root.find(f".//{{{self.DRAWING_NS}}}bodyPr")

            if body_pr is not None:
                body_pr.set("anchor", "ctr")

            if ring_center is not None:

                anchor_elm = root.find(f"{{{self.CHART_DRAWING_NS}}}relSizeAnchor")

                if anchor_elm is not None:

                    from_elm = anchor_elm.find(f"{{{self.CHART_DRAWING_NS}}}from")
                    to_elm = anchor_elm.find(f"{{{self.CHART_DRAWING_NS}}}to")

                    from_x_elm = from_elm.find(f"{{{self.CHART_DRAWING_NS}}}x")
                    from_y_elm = from_elm.find(f"{{{self.CHART_DRAWING_NS}}}y")
                    to_x_elm = to_elm.find(f"{{{self.CHART_DRAWING_NS}}}x")
                    to_y_elm = to_elm.find(f"{{{self.CHART_DRAWING_NS}}}y")

                    half_w = (float(to_x_elm.text) - float(from_x_elm.text)) / 2
                    half_h = (float(to_y_elm.text) - float(from_y_elm.text)) / 2

                    cx, cy = ring_center

                    from_x_elm.text = str(cx - half_w)
                    from_y_elm.text = str(cy - half_h)
                    to_x_elm.text = str(cx + half_w)
                    to_y_elm.text = str(cy + half_h)

            user_shapes_part.blob = serialize_part_xml(root)

    # ---------------------------------------------------------
    # Replace multiple text objects
    # ---------------------------------------------------------

    def replace_objects(self, replacements):

        print()

        print("=" * 60)
        print("UPDATING POWERPOINT OBJECTS")
        print("=" * 60)

        for key, value in replacements.items():

            if value is None:
                continue

            # A dict value carries a text payload plus targeting hints,
            # for objects whose name is ambiguous (occurrence) or where
            # the text belongs in a specific paragraph, not paragraph 0.
            # The dict key itself only needs to be unique in this dict --
            # it doesn't have to be the literal PowerPoint object name.
            if isinstance(value, dict) and "text" in value:

                object_name = value.get("object", key)
                occurrence = value.get("occurrence", 0)
                paragraph_index = value.get("paragraph_index", 0)
                group_child_index = value.get("group_child_index")
                run_index = value.get("run_index")
                payload = value["text"]

            else:

                object_name = key
                occurrence = 0
                paragraph_index = 0
                group_child_index = None
                run_index = None
                payload = value

            shape = self.find_shape(object_name, occurrence=occurrence)

            if shape is None:
                print(f"[NOT FOUND] {object_name}")
                continue

            if hasattr(shape, "has_chart") and shape.has_chart:

                self.replace_chart(
                    object_name,
                    payload
                )

            elif hasattr(shape, "has_table") and shape.has_table:

                self.replace_table(
                    object_name,
                    payload
                )

            elif isinstance(payload, list):

                self.replace_list(
                    object_name,
                    payload
                )

            else:

                self.replace_text(
                    object_name,
                    payload,
                    occurrence=occurrence,
                    paragraph_index=paragraph_index,
                    group_child_index=group_child_index,
                    run_index=run_index
                )

    # ---------------------------------------------------------
    # The template styles its "Status" column header two different
    # ways on two different slides -- navy (matching every other table
    # header in the deck) on the QoQ Metrics slide, but the theme's
    # orange accent color on the Optimization Highlights slide. Force
    # the orange one to match so the deck reads as one consistent
    # design rather than the template's own inconsistency.
    # ---------------------------------------------------------

    def _style_status_headers(self):

        navy = "262E81"

        def recolor(shapes):

            for shape in shapes:

                if (
                    shape.name == "CARD_StatusHeader"
                    and shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
                ):
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = RGBColor.from_string(navy)

                if hasattr(shape, "shapes"):
                    recolor(shape.shapes)

        for slide in self.prs.slides:
            recolor(slide.shapes)

    # ---------------------------------------------------------
    # The template's Top Intent table only has 2 columns
    # (Company, Score) -- relabel the first header to "Account" (the
    # deck talks about accounts throughout, and this is the same entity
    # as the Top Engaged Accounts table beside it) and grow it to 3
    # (Account, Topic, Score) so the real "what topic are they surging
    # on" detail has somewhere to go.
    # ---------------------------------------------------------

    def _prepare_intent_companies_table(self):

        table_shape = self.find_shape("Table_TopIntentCompanies")

        if table_shape is None or not table_shape.has_table:
            print("[SKIPPED] _prepare_intent_companies_table (table not found)")
            return

        table = table_shape.table

        # Done before the column-count check below so the relabel still
        # happens even when the table already has its 3rd column.
        header_tf = table.rows[0].cells[0].text_frame

        if header_tf.paragraphs and header_tf.paragraphs[0].runs:
            header_tf.paragraphs[0].runs[0].text = "Account"
        else:
            header_tf.text = "Account"

        if len(table.columns) >= 3:
            return

        original_widths = [col.width for col in table.columns]

        add_table_column(
            table,
            insert_at=1,
            header_text="Topic",
            clone_from_col_index=0
        )

        total_width = sum(original_widths)
        score_width = original_widths[-1]
        account_width = round(total_width * 0.40)
        topic_width = total_width - account_width - score_width

        set_column_widths(table, [account_width, topic_width, score_width])

        print("[RESIZED TABLE] Table_TopIntentCompanies -> added Topic column")

    # ---------------------------------------------------------
    # Merge the two period boxes on the Campaign Snapshot slide
    # into one wide box, for modes where both slots describe the
    # same window (e.g. Full Campaign) and showing it twice would
    # just be a duplicate.
    # ---------------------------------------------------------

    def merge_period_boxes(self):

        box = self.find_shape("PERIOD_Q1")
        background = self.find_shape("Rectangle 3")
        spare_box = self.find_shape("PERIOD_Q2")
        spare_background = self.find_shape("Rectangle 4")

        if not all([box, background, spare_box, spare_background]):
            print("[SKIPPED] merge_period_boxes (objects not found)")
            return False

        full_width = box.width + spare_box.width

        box.width = full_width
        background.width = full_width

        # Delete the now-redundant box/background outright rather than
        # parking them off the visible slide. Off-canvas shapes are not
        # harmless: PowerPoint sizes the editing pane's scrollable canvas
        # to enclose every shape, including ones outside the slide, so a
        # shape sitting ~23in to the left adds a horizontal scrollbar and
        # renders the slide flush against the right edge of the viewport
        # -- which reads as the slide's content being "stuck to the right
        # side" even though the content itself is correctly centered.
        for shape in (spare_box, spare_background):
            shape._element.getparent().remove(shape._element)

        print("[MERGED] PERIOD_Q1 + PERIOD_Q2 into a single box")

        # The template's own KPI card row sits noticeably closer to the
        # left edge of the merged bar than the right (a bigger gap on
        # the right than the left) -- present in the pristine template
        # itself, not something merging the boxes caused, but it reads
        # as the row being "stuck to one side" under the now-full-width
        # bar. Recenter the row (as a rigid block, same card widths and
        # spacing) within that same span.
        card_names = (
            "CARD_AssetsUsed", "CARD_JobTitles",
            "CARD_Country", "CARD_WeeksInMarket",
        )

        cards = [self.find_shape(name) for name in card_names]
        cards = [card for card in cards if card is not None]

        if cards:

            row_left = min(card.left for card in cards)
            row_right = max(card.left + card.width for card in cards)

            target_left = box.left + (box.width - (row_right - row_left)) / 2
            shift = round(target_left - row_left)

            for card in cards:
                card.left += shift

        return True

    # ---------------------------------------------------------
    # Reshape the deck's slide count to match how many performance
    # periods the selected analysis mode produced, before the fill
    # pass (replace_objects) runs.
    # ---------------------------------------------------------

    def _slide_index_by_shape(self, object_name):

        for index, slide in enumerate(self.prs.slides):
            for shape in slide.shapes:
                if shape.name == object_name:
                    return index

        return None

    @staticmethod
    def _rename_slot_shapes(slide, old_token, new_token):

        def walk(shapes):

            for shape in shapes:

                if old_token in shape.name:
                    shape.name = shape.name.replace(old_token, new_token)

                if hasattr(shape, "shapes"):
                    walk(shape.shapes)

        walk(slide.shapes)

    def reshape_for_periods(self, slots, comparison_slots):

        comparison_index = self._slide_index_by_shape("SECTION_Comparison")

        if comparison_index is None:
            print("[SKIPPED] reshape_for_periods (SECTION_Comparison not found)")
            return

        # In the template, the comparison divider is always preceded by
        # exactly two divider+detail pairs (Q1, then Q2).
        divider_index = comparison_index - 4
        detail_index = comparison_index - 3

        if len(slots) == 1:

            # Full Campaign: only one period -- the Q2 pair doesn't apply.
            delete_slide(self.prs, comparison_index - 1)
            delete_slide(self.prs, comparison_index - 2)

            print("[REMOVED] Q2 divider+detail pair (Full Campaign)")

        elif len(slots) > 2:

            # Month over Month with 3+ periods: clone the Q1 pair once
            # per extra slot, inserting each pair immediately before the
            # (still-unmoved) comparison section.
            cursor = comparison_index

            for slot in slots[2:]:

                new_divider = duplicate_slide(self.prs, divider_index, cursor)
                self._rename_slot_shapes(new_divider, "Q1", slot["slot"])
                cursor += 1

                new_detail = duplicate_slide(self.prs, detail_index, cursor)
                self._rename_slot_shapes(new_detail, "Q1", slot["slot"])
                cursor += 1

            print(f"[ADDED] {len(slots) - 2} extra divider+detail pair(s)")

        # Custom mode is the only case where the per-period detail
        # collapses to one slot while the comparison genuinely has two
        # (the single selected range, auto-bisected) -- that structural
        # signature, not the mode name, is what tells
        # _prepare_comparison_table() to relabel the period headers
        # even though there's no column growth to do.
        single_period_comparison = len(slots) == 1 and len(comparison_slots) == 2

        self._prepare_comparison_table(comparison_slots, single_period_comparison)

    @staticmethod
    def _abbreviate_period_label(label):

        """Shortens a full month name (or a "Month - Month" bucket
        range) to its 3-letter form so it fits on one line in a narrow
        comparison-table column instead of wrapping ("Septemb/er").
        Labels that aren't month names (quarter labels, "Period 1",
        custom-range labels) pass through unchanged."""

        if label in calendar.month_name[1:]:
            return label[:3]

        if " - " in label:

            parts = label.split(" - ")

            return "-".join(
                part[:3] if part in calendar.month_name[1:] else part
                for part in parts
            )

        return label

    def _prepare_comparison_table(self, comparison_slots, single_period_comparison=False):

        if len(comparison_slots) < 2:
            return

        if len(comparison_slots) == 2 and not single_period_comparison:
            return

        table_shape = self.find_shape("Table_QoQMetrics")

        if table_shape is None or not table_shape.has_table:
            print("[SKIPPED] _prepare_comparison_table (table not found)")
            return

        table = table_shape.table

        if single_period_comparison:

            # Exactly 2 comparison_slots, no column growth needed --
            # just swap the template's static "Q1"/"Q2" header text for
            # the real period labels ("Period 1"/"Period 2"). The
            # Change/% Change columns and their data stay exactly as
            # the template/QoQ Comparison table already have them.
            for i, slot in enumerate(comparison_slots):

                header_tf = table.rows[0].cells[i + 1].text_frame
                text = self._abbreviate_period_label(slot["label"])

                if header_tf.paragraphs and header_tf.paragraphs[0].runs:
                    header_tf.paragraphs[0].runs[0].text = text
                else:
                    header_tf.text = text

            print("[RELABELED TABLE] Table_QoQMetrics headers -> " +
                  ", ".join(slot["label"] for slot in comparison_slots))

            return

        # Capture the table's original layout before growing it --
        # add_table_column() clones the last column's width for each new
        # one, so without redistributing afterward, N extra columns just
        # make the table N times wider than the slide (pushing off the
        # right edge and overlapping whatever sits next to it), instead
        # of staying the same total width with narrower columns.
        original_total_width = sum(col.width for col in table.columns)
        original_metric_width = table.columns[0].width

        headers = (
            ["Metric"]
            + [self._abbreviate_period_label(slot["label"]) for slot in comparison_slots]
            + ["% Change"]
        )

        while len(table.columns) < len(headers):

            add_table_column(
                table,
                insert_at=len(table.columns),
                header_text="",
                clone_from_col_index=len(table.columns) - 1
            )

        for c in range(len(table.columns)):

            text = headers[c] if c < len(headers) else ""

            header_tf = table.rows[0].cells[c].text_frame

            if header_tf.paragraphs and header_tf.paragraphs[0].runs:
                header_tf.paragraphs[0].runs[0].text = text
            else:
                header_tf.text = text

            if c >= len(headers):

                for r in range(1, len(table.rows)):

                    cell_tf = table.rows[r].cells[c].text_frame

                    if cell_tf.paragraphs and cell_tf.paragraphs[0].runs:
                        cell_tf.paragraphs[0].runs[0].text = ""

        # Redistribute the table's original total width across however
        # many columns it now has -- Metric keeps a fixed share (it
        # doesn't grow with month count), the rest split evenly, so the
        # table stays within the slide instead of growing wider with
        # every extra month/quarter.
        other_count = len(table.columns) - 1

        if other_count > 0:

            other_width = round(
                (original_total_width - original_metric_width) / other_count
            )

            set_column_widths(
                table, [original_metric_width] + [other_width] * other_count
            )

        # Many columns leaves little width for each header/value -- scale
        # the font down as the count grows so text still fits instead of
        # wrapping or overflowing its cell.
        header_size = Pt(14) if len(table.columns) <= 5 else Pt(11) if len(table.columns) <= 7 else Pt(9)

        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = header_size

        print(
            f"[RESIZED TABLE] Table_QoQMetrics -> "
            f"{len(headers)} of {len(table.columns)} columns used"
        )

    # ---------------------------------------------------------
    # Save Presentation
    # ---------------------------------------------------------

    def save(self):

        output = Path("output") / "Generated_QBR.pptx"

        self.prs.save(output)

        print()
        print("=" * 60)
        print("POWERPOINT CREATED")
        print(output)
        print("=" * 60)

        return output

    # ---------------------------------------------------------
    # Remove a whole slide (identified by one of its own shapes) when
    # the report it depends on had no real data at all -- rather than
    # leaving the template's own static example content on screen
    # looking like real data for whatever client this deck is for.
    # ---------------------------------------------------------

    # ---------------------------------------------------------
    # Geography section (divider + market cards + market KPIs)
    #
    # The two content slides reuse shape names that already exist earlier in
    # the deck. The KPI slide was built from the Q2 performance slide and
    # kept its names -- PERIOD_Q2, CARD_Q2_TotalLeads, TITLE_Q2Performance --
    # and the market-card slide's shapes carry PowerPoint's own defaults
    # ("Text 4", "Shape 6"), which say nothing about what they hold.
    #
    # find_shape() resolves by name and takes the first match, so left as they
    # are the geography KPI cards would either never be filled (the Q2 slide
    # matches first) or, in Full Campaign mode where the Q2 pair is deleted,
    # be filled with a period's figures under a geography heading. Both slides
    # are renamed into their own GEO_* namespace immediately after load,
    # before anything else looks at the deck.
    #
    # The rename happens in the loaded copy only -- the template file on disk
    # keeps whatever names its author gave it.
    # ---------------------------------------------------------

    # The section is found by what its market slide is made of, not by a shape
    # name. The divider's name (SECTION_AudienceInterest) is already used by
    # the insights divider earlier in the deck, so a name lookup finds that one
    # and renames -- then deletes -- the topic and funnel slides instead. The
    # market slide is unmistakable by construction: it is the only slide built
    # from freeform markers, one per market.
    GEOGRAPHY_MIN_MARKERS = 2

    GEOGRAPHY_KPI_RENAMES = {
        "PERIOD_Q2": "GEO_Period",
        "TITLE_Q2Performance": "TITLE_GeographyPipeline",
        "AI_Q2_TopAsset": "AI_GeographyInsight",
        "CARD_Q2_TotalLeads": "GEO_CARD_1",
        "CARD_Q2_Accounts": "GEO_CARD_2",
        "CARD_Q2_Assets": "GEO_CARD_3",
        "CARD_Q2_JobTitles": "GEO_CARD_4",
        "CARD_Q2_Country": "GEO_CARD_5",
    }

    def namespace_geography_slides(self):

        """
        Renames the two geography content slides into the GEO_* namespace and
        reports how many market cards the template actually lays out.

        Returns the market-card count, or 0 when the section isn't in the
        template at all -- a template without these slides simply produces a
        deck without them rather than failing.
        """

        markets_index = self._market_slide_index()

        if markets_index is None:
            print("[SKIPPED] geography section (market slide not found)")
            return 0

        # The divider introduces the two slides around it, which is the only
        # relationship between them the template records -- nothing on the
        # divider names either content slide.
        kpi_index = markets_index + 1

        if kpi_index >= len(self.prs.slides):
            print("[SKIPPED] geography section (KPI slide missing)")
            return 0

        card_count = self._namespace_market_slide(self.prs.slides[markets_index])

        renamed = 0

        def walk(shapes):

            nonlocal renamed

            for shape in shapes:

                if shape.name in self.GEOGRAPHY_KPI_RENAMES:
                    shape.name = self.GEOGRAPHY_KPI_RENAMES[shape.name]
                    renamed += 1

                if hasattr(shape, "shapes"):
                    walk(shape.shapes)

        walk(self.prs.slides[kpi_index].shapes)

        print(
            f"[NAMESPACED] geography slides {markets_index + 1}-{kpi_index + 1} "
            f"({card_count} market cards, {renamed} KPI shapes)"
        )

        return card_count

    def _market_slide_index(self):

        """The slide built from freeform market markers, or None."""

        for index, slide in enumerate(self.prs.slides):

            markers = sum(
                1 for shape in slide.shapes
                if shape.shape_type == MSO_SHAPE_TYPE.FREEFORM
            )

            if markers >= self.GEOGRAPHY_MIN_MARKERS:
                return index

        return None

    @staticmethod
    def _namespace_market_slide(slide):

        """
        Names the market-card slide's shapes by where they sit, not by what
        they are called.

        Every card is a freeform marker with a name box and a detail box drawn
        inside it. The boxes have no names worth matching on, but their
        geometry is unambiguous: each sits within its own marker's bounds, and
        within a marker the name is above the detail. Reading the layout means
        the wiring survives the author adding, removing or re-styling cards --
        which a fixed "Text 4 is the first country" map would not.
        """

        markers = sorted(
            (
                shape for shape in slide.shapes
                if shape.shape_type == MSO_SHAPE_TYPE.FREEFORM
            ),
            key=lambda shape: shape.left,
        )

        def centre(shape):
            return (shape.left + shape.width / 2, shape.top + shape.height / 2)

        claimed = set()

        for index, marker in enumerate(markers, start=1):

            marker.name = f"GEO_Market{index}_Shape"

            inside = []

            for shape in slide.shapes:

                if shape is marker or not shape.has_text_frame:
                    continue

                if not shape.text_frame.text.strip():
                    continue

                x, y = centre(shape)

                if (
                    marker.left <= x <= marker.left + marker.width
                    and marker.top <= y <= marker.top + marker.height
                ):
                    inside.append(shape)

            inside.sort(key=lambda shape: shape.top)

            for label, shape in zip(("Name", "Detail"), inside):
                shape.name = f"GEO_Market{index}_{label}"
                claimed.add(id(shape))

        # Whatever full-width text box is left over is the slide's lead-in
        # line. The title is excluded by name rather than by width -- it is
        # very nearly as wide as the intro, so size alone would not separate
        # them reliably.
        leftovers = [
            shape for shape in slide.shapes
            if shape.has_text_frame
            and shape.text_frame.text.strip()
            and id(shape) not in claimed
            and not shape.name.startswith("GEO_Market")
        ]

        for shape in leftovers:

            if shape.name.startswith("TITLE_"):
                shape.name = "TITLE_GeographyMarkets"

        intro = [
            shape for shape in leftovers
            if not shape.name.startswith("TITLE_")
        ]

        if intro:
            max(intro, key=lambda shape: shape.width).name = "AI_GeographyIntro"

        # The rule the cards are threaded onto. It spans exactly the block of
        # cards, so it has to be resized whenever the block is -- otherwise it
        # runs on past the last card into empty slide. It is the only wide
        # shape here carrying no text.
        axis = [
            shape for shape in slide.shapes
            if shape.shape_type != MSO_SHAPE_TYPE.FREEFORM
            and not shape.name.startswith(("GEO_Market", "TITLE_", "AI_"))
            and not (shape.has_text_frame and shape.text_frame.text.strip())
        ]

        if axis:
            max(axis, key=lambda shape: shape.width).name = "GEO_MarketAxis"

        return len(markers)

    @staticmethod
    def _market_cards(slide):

        """{position: [shapes]} for each market card on the slide."""

        cards = {}

        for shape in slide.shapes:

            match = re.match(r"GEO_Market(\d+)_", shape.name)

            if match:
                cards.setdefault(int(match.group(1)), []).append(shape)

        return cards

    @staticmethod
    def _card_marker(shapes):

        """The freeform a card is drawn on -- the shape whose position the
        card's labels are placed relative to."""

        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.FREEFORM:
                return shape

        return max(shapes, key=lambda shape: shape.width)

    def trim_market_cards(self, keep):

        """
        Removes the market cards the campaign has no market for, then
        recentres the ones that remain.

        The template draws a fixed number of cards; a campaign selling into
        three countries fills three of them. Leaving the rest showing the
        template's own example markets would put another company's countries
        on the slide, so the surplus is deleted outright -- but deleting alone
        leaves the survivors bunched against the left edge with a bay of empty
        slide to their right, and the rule they sit on still running the full
        width into nothing.

        The cards keep the spacing the template gave them. They are drawn as a
        near-continuous chain -- a quarter-inch between two-inch cards -- and
        that chain is the design; spreading three cards over the width the
        template used for five would open ten times the gap and lose it. So
        the block keeps its pitch, is recentred on the same axis the full row
        was centred on, and the rule is resized to span exactly the cards that
        are left, which is the relationship the template already holds.
        """

        index = self._market_slide_index()

        if index is None:
            return

        slide = self.prs.slides[index]

        cards = self._market_cards(slide)

        if not cards:
            return

        positions = sorted(cards)

        markers = {
            position: self._card_marker(shapes)
            for position, shapes in cards.items()
        }

        # Measured before anything is deleted: the pitch and the centre both
        # describe the row the template drew, which is what the trimmed row
        # has to go on matching.
        pitch = (
            round(
                (markers[positions[-1]].left - markers[positions[0]].left)
                / (len(positions) - 1)
            )
            if len(positions) > 1 else 0
        )

        block_left = markers[positions[0]].left

        block_right = markers[positions[-1]].left + markers[positions[-1]].width

        centre = (block_left + block_right) / 2

        removed = 0

        for position in positions:

            if position > keep:

                for shape in cards.pop(position):
                    shape._element.getparent().remove(shape._element)
                    removed += 1

        if not removed:
            return

        print(f"[REMOVED] {removed} unused market card shape(s)")

        self._recentre_market_cards(cards, markers, pitch, centre)

    def _recentre_market_cards(self, cards, markers, pitch, centre):

        kept = sorted(cards)

        if not kept or not pitch:
            return

        width = markers[kept[0]].width

        block = (len(kept) - 1) * pitch + width

        left = centre - block / 2

        for slot, position in enumerate(kept):

            shift = round(left + slot * pitch - markers[position].left)

            if not shift:
                continue

            # Every shape of the card moves together -- the marker and the two
            # labels drawn inside it are separate shapes, not a group.
            for shape in cards[position]:
                shape.left += shift

        axis = self.find_shape("GEO_MarketAxis")

        if axis is not None:
            axis.left = round(left)
            axis.width = round(block)

        print(
            f"[RECENTRED] {len(kept)} market cards across "
            f"{Emu(round(block)).inches:.2f}in"
        )

    def remove_geography_section(self):

        """
        Drops all three geography slides.

        A campaign delivering into a single country has nothing to break down
        by market: the cards would be one filled and the rest deleted, and the
        concentration KPIs would all read 100%. The country count on the
        Campaign Snapshot already says everything there is to say.
        """

        markets_index = self._market_slide_index()

        if markets_index is None:
            return

        doomed = [markets_index, markets_index + 1]

        # The divider before the market slide goes too, but only once
        # confirmed to be a divider -- a section heading and nothing else.
        # Deleting by position alone would take a content slide with it if the
        # section were ever reordered.
        if markets_index:

            previous = self.prs.slides[markets_index - 1]

            if any(
                shape.name.startswith("SECTION_")
                for shape in previous.shapes
            ):
                doomed.append(markets_index - 1)

        for index in sorted(doomed, reverse=True):
            if index < len(self.prs.slides):
                delete_slide(self.prs, index)

        print(
            f"[REMOVED] geography section ({len(doomed)} slides, "
            "fewer than 2 markets)"
        )

    def fit_footnote(self, object_name, reference_names, min_height, gap=45720):

        """
        Keeps a one-line footnote box clear of the content above it and gives
        it room for the text it actually received.

        Two template traits make these boxes collide with their own slide.
        They are sized for a single short line (AI_TopAccountsFooter is 0.17in
        tall), so real commentary wraps to two or three lines and spills out of
        the box -- upward, over the tables above it. And some sit slightly
        inside the shape above them to begin with.

        Rather than hardcode positions per slide, this measures: it drops the
        box below everything in `reference_names`, grows it to `min_height`,
        and enables word wrap. Any campaign, any client, any text length --
        the box moves to fit rather than the text being cut to fit.
        """

        shape = self.find_shape(object_name)

        if shape is None:
            return

        references = [
            self.find_shape(name) for name in reference_names
        ]

        bottoms = [
            r.top + r.height
            for r in references
            if r is not None and r.top is not None and r.height is not None
        ]

        if not bottoms:
            return

        original_top, original_height = shape.top, shape.height

        if shape.height < min_height:
            shape.height = min_height

        # Anchored to the bottom of the slide rather than to the bottom of the
        # content above it. The tables above these footnotes contain client
        # names, and PowerPoint grows a row at render time to fit a name that
        # wraps -- so shape.height under-reports how tall the table really is,
        # and "below the table" computed from it can still land inside the
        # rendered table. Sitting just above the footer band is the position
        # that holds for any client's data.
        floor_top = (
            self.prs.slide_height - FOOTER_BAND_HEIGHT - shape.height
        )

        target_top = max(max(bottoms) + gap, floor_top)

        if shape.top < target_top:
            shape.top = target_top

        # Without this a long line runs off the side of the slide instead of
        # wrapping into the height just made available for it.
        if shape.has_text_frame:
            shape.text_frame.word_wrap = True

        # Never push the box off the bottom of the slide; if the space simply
        # isn't there, keep it on-slide and let the smaller height apply.
        slide_height = self.prs.slide_height

        if shape.top + shape.height > slide_height:
            shape.height = max(original_height, slide_height - shape.top)

        if (shape.top, shape.height) != (original_top, original_height):
            print(
                f"[REFLOWED] {object_name} top {original_top} -> {shape.top}, "
                f"height {original_height} -> {shape.height}"
            )

    def clear_comparison_summary_overlap(self):

        """
        Moves the comparison slide's one-line summary clear of the chart above
        it.

        In the template AI_ComparisonSummary starts at 5756105 EMU while
        CHART_Q1Q2Comparison runs to 5938985 -- a 0.2in overlap, so the summary
        sits on top of the chart's category axis labels. Nudging the text box
        down is the code-side fix for a template geometry problem; the template
        itself is never edited.
        """

        summary = self.find_shape("AI_ComparisonSummary")
        chart = self.find_shape("CHART_Q1Q2Comparison")

        if summary is None or chart is None:
            return

        gap = Emu(45720)  # 0.05in

        chart_bottom = chart.top + chart.height

        if summary.top >= chart_bottom + gap:
            return

        original_top = summary.top

        summary.top = chart_bottom + gap

        print(
            f"[MOVED] AI_ComparisonSummary {original_top} -> {summary.top} "
            f"(clear of chart bottom {chart_bottom})"
        )

    def apply_commentary_body_font(self):

        """Resizes the body paragraphs of the fixed-size commentary cards (see
        COMMENTARY_BODY_FONT). Paragraph 0 is each card's own bold label and is
        left at the template's size."""

        for object_name, size in COMMENTARY_BODY_FONT.items():

            shape = self.find_shape(object_name)

            if shape is None or not shape.has_text_frame:
                continue

            for paragraph in shape.text_frame.paragraphs[1:]:
                for run in paragraph.runs:
                    run.font.size = size

            print(f"[RESIZED BODY] {object_name} -> {size.pt:g}pt")

    # Rough metrics for estimating how tall wrapped text will render. The
    # deck uses one humanist sans throughout, whose average glyph runs a little
    # over half its point size; 1.2 is the line height PowerPoint applies at
    # single spacing. Both are approximations, which is why shrink_to_fit()
    # keeps a margin in hand rather than filling the box to the last point.
    CHAR_WIDTH_RATIO = 0.52
    LINE_HEIGHT_RATIO = 1.2
    FIT_SAFETY_MARGIN = 0.92

    def shrink_to_fit(self, object_name, min_scale=0.7):

        """
        Scales a text box's font down until its content fits the box.

        Needed wherever the number of entries isn't fixed. Key Learnings is
        drawn for five, and a campaign with a market breakdown gets a sixth --
        which at the template's own font size runs off the bottom of the
        slide. Rather than sizing for the worst case and leaving five-entry
        decks looking sparse, the text keeps its designed size whenever it
        fits and only gives ground when it doesn't.

        PowerPoint's own shrink-on-overflow is not an option here: it is
        applied by the application when the box is edited, so a deck rendered
        or exported without ever being opened keeps overflowing.
        """

        shape = self.find_shape(object_name)

        if shape is None or not shape.has_text_frame:
            return

        frame = shape.text_frame

        width = shape.width - (frame.margin_left or 0) - (frame.margin_right or 0)
        height = shape.height - (frame.margin_top or 0) - (frame.margin_bottom or 0)

        usable_width = Emu(width).pt
        usable_height = Emu(height).pt * self.FIT_SAFETY_MARGIN

        # (text, point size, fixed spacing) per paragraph. Spacing is left
        # unscaled -- it is set on the paragraph, not the run, and shrinking
        # the font doesn't change it.
        lines = []

        for paragraph in frame.paragraphs:

            size = next(
                (run.font.size.pt for run in paragraph.runs if run.font.size),
                None,
            )

            spacing = sum(
                value.pt for value in
                (paragraph.space_before, paragraph.space_after)
                if value is not None
            )

            lines.append((paragraph.text, size, spacing))

        # Some runs carry no explicit size and inherit one -- the bold titles
        # in Key Learnings are written that way. They still take up space and
        # still have to shrink, so they borrow the size of a sized run in the
        # same box, which in these template boxes is the size they render at.
        default_size = next(
            (size for _, size, _ in lines if size), 15.0
        )

        lines = [
            (text, size or default_size, spacing)
            for text, size, spacing in lines
        ]

        def height_at(scale):

            total = 0.0

            for text, size, spacing in lines:

                size *= scale

                per_line = max(1, usable_width / (size * self.CHAR_WIDTH_RATIO))

                wrapped = max(1, math.ceil(len(text) / per_line))

                total += wrapped * size * self.LINE_HEIGHT_RATIO + spacing

            return total

        scale = 1.0

        while scale > min_scale and height_at(scale) > usable_height:
            scale -= 0.05

        if scale >= 1.0:
            return

        for paragraph in frame.paragraphs:
            for run in paragraph.runs:
                current = run.font.size.pt if run.font.size else default_size
                run.font.size = Pt(round(current * scale, 1))

        print(f"[SHRUNK] {object_name} to {scale:.0%} to fit its box")

    def _delete_slides_by_shape(self, shape_names):

        for name in shape_names:

            index = self._slide_index_by_shape(name)

            if index is not None:
                delete_slide(self.prs, index)
                print(f"[REMOVED] slide containing {name} (no source data)")

    # ---------------------------------------------------------
    # Create PPT
    # ---------------------------------------------------------

    def create(
        self,
        replacements=None,
        merge_period_boxes=False,
        periods=None,
        empty_data_slides=None,
        market_count=0
    ):

        self.load()

        # First of all, before any other pass reads a shape by name: the
        # geography slides share several names with the Q2 performance slide,
        # and until they are renamed a lookup for either can land on the wrong
        # one.
        market_slots = self.namespace_geography_slides()

        if periods:

            self.reshape_for_periods(
                periods.get("slots", []),
                periods.get("comparison_slots", [])
            )

        if market_slots and market_count > 1:

            self.trim_market_cards(min(market_count, market_slots))

            # Key Learnings and the recommendations list each gain one entry
            # for geography. Grown here, before the fill pass, so the extra
            # paragraphs exist by the time their text is written.
            self.clone_paragraphs("AI_KeyLearnings", (-2, -1))
            self.clone_paragraphs("AI_H2Recommendations", (-1,))

        else:
            self.remove_geography_section()

        if empty_data_slides:
            self._delete_slides_by_shape(empty_data_slides)

        self._prepare_intent_companies_table()

        if replacements:

            self.replace_objects(replacements)

            # After the text is in, not before -- replace_text() rewrites the
            # run that carries the size.
            self.apply_commentary_body_font()

            self.clear_comparison_summary_overlap()

            # Variable-length lists, after their text is in: how far they
            # overflow depends on what the campaign produced, so it can only
            # be measured once written.
            self.shrink_to_fit("AI_KeyLearnings")
            self.shrink_to_fit("AI_H2Recommendations")

            # Footnote boxes the template sized for one short line, which real
            # commentary overflows into the content above them.
            self.fit_footnote(
                "AI_TopAccountsFooter",
                ["Table_TopEngagedAccounts", "Table_TopIntentCompanies"],
                min_height=Emu(600000),
            )

            self.fit_footnote(
                "AI_OptimizationFooter",
                ["Table_OptimizationHighlights"],
                min_height=Emu(500000),
            )

        if merge_period_boxes:

            self.merge_period_boxes()

        self._style_status_headers()

        self.save()
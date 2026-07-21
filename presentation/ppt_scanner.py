from pptx import Presentation


class PPTScanner:

    def __init__(self, ppt_path):

        self.prs = Presentation(ppt_path)

    def scan(self):

        with open("Objects.txt", "w", encoding="utf-8") as f:

            for slide_no, slide in enumerate(self.prs.slides, start=1):

                f.write(f"\n========== SLIDE {slide_no} ==========\n")

                for shape in slide.shapes:

                    shape_type = type(shape).__name__

                    f.write(f"{shape.name} [{shape_type}]\n")

                    if hasattr(shape, "has_table") and shape.has_table:

                        f.write("   -> TABLE\n")

                    elif hasattr(shape, "has_chart") and shape.has_chart:

                        chart = shape.chart

                        f.write("   -> CHART\n")

                        try:
                            f.write(f"   -> TYPE: {chart.chart_type}\n")
                        except Exception:
                            pass

                    elif hasattr(shape, "has_text_frame") and shape.has_text_frame:

                        text = shape.text.replace("\n", " ")

                        if len(text) > 120:
                            text = text[:120] + "..."

                        f.write(f"   -> {text}\n")
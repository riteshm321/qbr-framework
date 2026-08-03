"""
Slide/Table Reshaping Operations

python-pptx has no built-in slide-duplication, slide-deletion, or
table-column-insertion API. These are the primitives the QBR needs to
adapt the template's fixed slide count to however many periods the
selected analysis mode produces (e.g. one divider+detail slide pair per
month for Month over Month).

Scoped deliberately to what this template actually uses: slides here only
ever relate to a slideLayout and, on the comparison slide, a native chart.
There are no images, hyperlinks, or notes slides in this deck.
"""

import copy

from pptx.opc.constants import CONTENT_TYPE as CT
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.package import Part
from pptx.oxml.ns import qn
from pptx.parts.chart import ChartPart
from pptx.parts.embeddedpackage import EmbeddedXlsxPart
from pptx.parts.slide import SlidePart

# Not exposed as a named RT constant in this python-pptx version.
CHART_STYLE = "http://schemas.microsoft.com/office/2011/relationships/chartStyle"


def _clone_generic_part(package, source_part, partname_template):

    """
    Clones a non-XML-parsed part (chartStyle/chartColorStyle load as plain
    `Part` objects, not `XmlPart`, in this python-pptx version) by copying
    its raw bytes rather than an lxml element.
    """

    new_partname = package.next_partname(partname_template)

    return Part(new_partname, source_part.content_type, package, source_part.blob)


def _clone_chart_part(package, source_chart_part):

    """
    Clones a chart part along with its own embedded workbook and its
    chartStyle/chartColorStyle parts.

    All three must be independently owned, not shared with the source
    chart -- confirmed empirically against real PowerPoint (not just from
    reading python-pptx's source): sharing chartStyle/chartColorStyle
    between two chart parts, even though neither is referenced by an
    inline r:id anywhere, causes PowerPoint to refuse to open the file
    outright (no repair prompt, just "could not open the file"). The
    embedded workbook must never be shared either, since
    `chart.replace_data()` mutates whatever part its rel currently
    resolves to -- a shared workbook would let editing one chart corrupt
    the other's backing data.
    """

    new_element = copy.deepcopy(source_chart_part._element)

    new_partname = package.next_partname(ChartPart.partname_template)

    new_chart_part = ChartPart(new_partname, CT.DML_CHART, package, new_element)

    for reltype, partname_template in (
        (CHART_STYLE, "/ppt/charts/style%d.xml"),
        (RT.CHART_COLOR_STYLE, "/ppt/charts/colors%d.xml"),
    ):

        try:
            source_related = source_chart_part.part_related_by(reltype)
        except KeyError:
            continue

        cloned_related = _clone_generic_part(package, source_related, partname_template)

        new_chart_part.relate_to(cloned_related, reltype)

    source_xlsx_part = source_chart_part.chart_workbook.xlsx_part

    cloned_xlsx_part = EmbeddedXlsxPart.new(source_xlsx_part.blob, package)

    new_chart_part.chart_workbook.xlsx_part = cloned_xlsx_part

    return new_chart_part


def _rewrite_inline_rid(root_element, attr_qname, old_rid, new_rid):

    for element in root_element.iter():

        if element.get(attr_qname) == old_rid:
            element.set(attr_qname, new_rid)


def duplicate_slide(prs, source_index, insert_before_index):

    """
    Deep-copies the slide at `source_index` and inserts the copy at
    `insert_before_index`. The slideLayout relationship is shared (it's
    ambient -- never referenced by an inline r:id); a chart, if present,
    is independently cloned via `_clone_chart_part`.

    Returns the new `Slide`. Caller is responsible for renaming the
    clone's shapes before running the fill pass, since `find_shape()`
    resolves by exact name and a same-named clone would be permanently
    unaddressable.
    """

    source_slide = prs.slides[source_index]
    source_part = source_slide.part

    new_element = copy.deepcopy(source_part._element)

    new_partname = prs.part.package.next_partname("/ppt/slides/slide%d.xml")

    new_slide_part = SlidePart(new_partname, CT.PML_SLIDE, prs.part.package, new_element)

    for rel in source_part.rels.values():

        if rel.reltype == RT.SLIDE_LAYOUT:

            new_slide_part.relate_to(rel.target_part, RT.SLIDE_LAYOUT)

        elif rel.reltype == RT.CHART:

            cloned_chart_part = _clone_chart_part(prs.part.package, rel.target_part)

            new_rid = new_slide_part.relate_to(cloned_chart_part, RT.CHART)

            _rewrite_inline_rid(new_element, qn("r:id"), rel.rId, new_rid)

        elif rel.reltype == RT.IMAGE:

            new_rid = new_slide_part.relate_to(rel.target_part, RT.IMAGE)

            _rewrite_inline_rid(new_element, qn("r:embed"), rel.rId, new_rid)

        else:

            raise NotImplementedError(
                f"duplicate_slide: unhandled relationship type {rel.reltype}"
            )

    slide_rid = prs.part.relate_to(new_slide_part, RT.SLIDE)

    sld_id_lst = prs.slides._sldIdLst

    new_sld_id = sld_id_lst.add_sldId(slide_rid)

    sld_id_lst.remove(new_sld_id)

    sld_id_lst.insert(insert_before_index, new_sld_id)

    for index, sld_id in enumerate(sld_id_lst.sldId_lst):

        if sld_id.rId == slide_rid:
            return prs.slides[index]

    raise RuntimeError("duplicate_slide: could not locate the newly inserted slide")


def delete_slide(prs, index):

    """Removes the slide at `index`. Order matters: the `<p:sldId>` must
    be removed before the relationship is dropped, since `drop_rel` only
    drops a relationship still referenced fewer than twice."""

    sld_id_lst = prs.slides._sldIdLst

    sld_id = sld_id_lst.sldId_lst[index]

    rid = sld_id.rId

    sld_id_lst.remove(sld_id)

    prs.part.drop_rel(rid)


def add_table_column(table, insert_at, header_text, clone_from_col_index=None):

    """
    Grows a native PPT table by one column, cloning an existing column's
    formatting (fill, font, borders) rather than inserting a blank one.
    Requires a uniform grid (no merged cells) -- true for every table in
    this template.
    """

    if clone_from_col_index is None:
        clone_from_col_index = len(table.columns) - 1

    tbl = table._tbl

    grid_col = copy.deepcopy(tbl.tblGrid.gridCol_lst[clone_from_col_index])

    tbl.tblGrid.insert(insert_at, grid_col)

    for row in tbl.findall(qn("a:tr")):

        cells = row.findall(qn("a:tc"))

        new_cell = copy.deepcopy(cells[clone_from_col_index])

        row.insert(insert_at, new_cell)

    header_cell = table.rows[0].cells[insert_at]

    tf = header_cell.text_frame

    if tf.paragraphs and tf.paragraphs[0].runs:
        tf.paragraphs[0].runs[0].text = str(header_text)
    else:
        tf.text = str(header_text)


def set_column_widths(table, widths):

    """Sets each column's width (EMU) directly, e.g. to redistribute a
    table's original total width evenly across a new column count."""

    for column, width in zip(table.columns, widths):
        column.width = width
